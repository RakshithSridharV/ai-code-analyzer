import collections.abc
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

def create_presentation():
    prs = Presentation()
    
    # We will use the blank layout (usually layout 6) to have full control
    blank_layout = prs.slide_layouts[6]

    def add_slide_with_header(title_text):
        slide = prs.slides.add_slide(blank_layout)
        
        # Add background header shape
        header_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Inches(1.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(41, 128, 185) # Blue
        header_shape.line.color.rgb = RGBColor(41, 128, 185)
        
        # Add title text
        txBox = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), prs.slide_width - Inches(1), Inches(0.8))
        tf = txBox.text_frame
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.text = title_text
        p.font.bold = True
        p.font.size = Pt(36)
        p.font.color.rgb = RGBColor(255, 255, 255)
        
        # Add footer
        footer_shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, prs.slide_height - Inches(0.5), prs.slide_width, Inches(0.5)
        )
        footer_shape.fill.solid()
        footer_shape.fill.fore_color.rgb = RGBColor(44, 62, 80) # Dark Navy
        footer_shape.line.color.rgb = RGBColor(44, 62, 80)
        
        footer_txBox = slide.shapes.add_textbox(Inches(0.5), prs.slide_height - Inches(0.5), prs.slide_width - Inches(1), Inches(0.5))
        f_tf = footer_txBox.text_frame
        f_tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        f_p = f_tf.paragraphs[0]
        f_p.text = "INDIRA INSTITUTE OF ENGINEERING & TECHNOLOGY | Dept. of CSE | 2022–2026"
        f_p.font.size = Pt(12)
        f_p.font.color.rgb = RGBColor(255, 255, 255)
        
        return slide

    # Slide 1: Title
    slide1 = prs.slides.add_slide(blank_layout)
    bg = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(44, 62, 80)
    bg.line.color.rgb = RGBColor(44, 62, 80)
    
    t_box = slide1.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(2))
    t_tf = t_box.text_frame
    t_p = t_tf.paragraphs[0]
    t_p.text = "ASTra"
    t_p.font.bold = True
    t_p.font.size = Pt(60)
    t_p.font.color.rgb = RGBColor(255, 255, 255)
    t_p.alignment = PP_ALIGN.CENTER
    
    sub_p = t_tf.add_paragraph()
    sub_p.text = "AST-Native Code Intelligence Platform"
    sub_p.font.size = Pt(28)
    sub_p.font.color.rgb = RGBColor(189, 195, 199)
    sub_p.alignment = PP_ALIGN.CENTER
    
    team_box = slide1.shapes.add_textbox(Inches(1), Inches(5), prs.slide_width - Inches(2), Inches(2))
    team_tf = team_box.text_frame
    team_p = team_tf.paragraphs[0]
    team_p.text = "Team: V. Rakshith Sridhar, N. Mohan, M. Sudarson\nGuide: Mr. P. Thanigesan"
    team_p.font.size = Pt(18)
    team_p.font.color.rgb = RGBColor(255, 255, 255)
    team_p.alignment = PP_ALIGN.CENTER
    
    # Slide 2: What is ASTra?
    slide2 = add_slide_with_header("What is ASTra?")
    content = [
        "A static code analysis platform that reads source code without running it.",
        "Infers Time & Space Complexity (O(n), O(n²), etc.).",
        "Detects inefficiencies, bugs, and bad patterns.",
        "Built completely from scratch using Python's ast module (no linter libraries borrowed).",
        "Includes an AI chat assistant and a machine-learning quality predictor."
    ]
    cb = slide2.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(4))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, point in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        
    # Slide 3: System Architecture
    slide3 = add_slide_with_header("System Architecture & Stack")
    content = [
        "Frontend: React 18 + Vite (Custom CSS glassmorphism, Recharts).",
        "Backend: Flask REST API (JWT Auth, CORS, Rate Limiting).",
        "Analysis Core: 10 Hand-written Engines (Python ast + Tree-sitter).",
        "AI Assistant: HuggingFace Inference API (Qwen 2.5 Coder 32B).",
        "Database: SQLite + SQLAlchemy.",
        "Infrastructure: Docker, GitHub Actions CI."
    ]
    cb = slide3.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(4))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, point in enumerate(content):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + point
        p.font.size = Pt(24)
        p.space_after = Pt(14)
        
    # Slide 4: 10 Analysis Engines
    slide4 = add_slide_with_header("The 10 Original Analysis Engines")
    engines = [
        "1. InferenceEngine: Time/space complexity from loop structure.",
        "2. RecursionClassifier: Classifies 6 recursion patterns.",
        "3. ExplanationBuilder: Deterministic plain-English explanations.",
        "4. DataFlowTracer: Tracks 5 data-flow inefficiency patterns.",
        "5. AntiPatternDetector: Catches 5 classic Python anti-patterns.",
        "6. CyclomaticAnalyzer: McCabe's complexity with risk labels.",
        "7. ConfidenceEstimator: Honest uncertainty bounds.",
        "8. DeadCodeDetector: Unused vars, imports, unreachable code.",
        "9. TypeInferencer: Static inference from literal assignments.",
        "10. HalsteadAnalyzer: Software science metrics (effort, bug estimate)."
    ]
    cb = slide4.shapes.add_textbox(Inches(0.5), Inches(1.5), prs.slide_width - Inches(1), Inches(5))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, eng in enumerate(engines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = eng
        p.font.size = Pt(20)
        p.space_after = Pt(6)

    # Slide 5: Multi-Language Support
    slide5 = add_slide_with_header("Multi-Language Support")
    langs = [
        "Python (Built-in ast): Full support — all 10 engines.",
        "JavaScript (Tree-sitter): Full complexity analysis, higher-order methods.",
        "Java (Tree-sitter): Full complexity, per-method breakdown.",
        "C (Tree-sitter): Full complexity, pointer/array allocation patterns.",
        "C++ (Tree-sitter): Full complexity, full STL container and template support."
    ]
    cb = slide5.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(4))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, lang in enumerate(langs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + lang
        p.font.size = Pt(24)
        p.space_after = Pt(14)

    # Slide 6: Testing & Validation
    slide6 = add_slide_with_header("Testing & Validation")
    tests = [
        "Unit Tests: pytest suite for all 10 engines + GitHub Actions CI.",
        "Cross-Language Tests: Consistency of complexity inference across parsers.",
        "Integration Tests: Full API endpoint testing & JWT auth flow validation.",
        "ML Model: Random Forest Classifier trained on Google MBPP dataset.",
        "Security: flask-talisman headers, bcrypt hashing, DOMPurify XSS protection."
    ]
    cb = slide6.shapes.add_textbox(Inches(1), Inches(2), prs.slide_width - Inches(2), Inches(4))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, test in enumerate(tests):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = "• " + test
        p.font.size = Pt(24)
        p.space_after = Pt(14)

    # Slide 7: Conclusion
    slide7 = add_slide_with_header("Conclusion & Future Work")
    conc = [
        "Achievements:",
        "  - 10 original analysis engines, zero borrowed logic.",
        "  - Deterministic Big-O inference with 97%+ confidence.",
        "  - Multi-language support and comprehensive platform built from scratch.",
        "Future Work:",
        "  - Space complexity inference (heap allocation tracking).",
        "  - Interprocedural analysis (cross-function complexity).",
        "  - IDE Plugin (VS Code extension) and support for Go, Rust."
    ]
    cb = slide7.shapes.add_textbox(Inches(1), Inches(1.5), prs.slide_width - Inches(2), Inches(5))
    tf = cb.text_frame
    tf.word_wrap = True
    for i, item in enumerate(conc):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.font.size = Pt(22)
        p.space_after = Pt(10)

    # Slide 8: Thank You
    slide8 = prs.slides.add_slide(blank_layout)
    bg = slide8.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(44, 62, 80)
    bg.line.color.rgb = RGBColor(44, 62, 80)
    
    t_box = slide8.shapes.add_textbox(Inches(1), Inches(3), prs.slide_width - Inches(2), Inches(2))
    t_tf = t_box.text_frame
    t_p = t_tf.paragraphs[0]
    t_p.text = "Thank You!"
    t_p.font.bold = True
    t_p.font.size = Pt(60)
    t_p.font.color.rgb = RGBColor(255, 255, 255)
    t_p.alignment = PP_ALIGN.CENTER
    
    sub_p = t_tf.add_paragraph()
    sub_p.text = "Questions?"
    sub_p.font.size = Pt(28)
    sub_p.font.color.rgb = RGBColor(189, 195, 199)
    sub_p.alignment = PP_ALIGN.CENTER

    prs.save("ASTra_Final_Review_Presentation.pptx")
    print("Presentation generated successfully!")

if __name__ == '__main__':
    create_presentation()
